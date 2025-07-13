# streamlit run gen-story.py --logger.level error
# streamlit run gen-story.py --server.fileWatcherType=none
# 
# File type detection uses built-in Python libraries and file signatures
# No additional dependencies required for Streamlit Cloud deployment
#
# Suppress Streamlit file watcher and set log level before any imports
import os
os.environ["PYTORCH_JIT"] = "0"
os.environ["PYTORCH_ENABLE_MPS_FALLBACK"] = "1"
os.environ["GOOGLE_API_USE_MTLS_ENDPOINT"] = "never"
os.environ["GRPC_DNS_RESOLVER"] = "native"
os.environ["CURL_CA_BUNDLE"] = ""
# Standard library imports
import asyncio
import html2text
import sys
import types
import queue
import tempfile
import io
import gc
import functools
import datetime
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor

# Third-party imports
import nest_asyncio
import streamlit as st
import markdown 
# from html2docx import html2docx # Only use function html2docx
from googleapiclient.discovery import build
from google.oauth2 import service_account
from googleapiclient.http import MediaIoBaseDownload, MediaIoBaseUpload
import google.generativeai as Ggenai
import json
import soundfile as sf
import numpy as np
from PIL import Image
import fitz  # PyMuPDF
import docx
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.shared import Pt
from docx.oxml.ns import qn
import pptx
from dotenv import load_dotenv
import aiofiles
from bs4 import BeautifulSoup
import torch
from googleapiclient.errors import HttpError
import time
import re
import socket
import ssl
import mimetypes
import markdownify
import subprocess

def update_progress(message):
    if 'progress_history' in st.session_state:
        st.session_state.progress_history.insert(0, message)
        st.session_state.progress_history.insert(1, '-'*39 + '\n')

def detect_file_type(file_data, file_meta=None):
    """Detect file type from content using file signatures and MIME types"""
    
    try:
        # First, try to get MIME type from Google Drive metadata
        if file_meta and 'mimeType' in file_meta:
            return file_meta['mimeType']
        
        # For local files or when metadata is not available, use file signatures
        if isinstance(file_data, str):  # Local file path
            # Use mimetypes for local files
            mime_type, _ = mimetypes.guess_type(file_data)
            return mime_type
        else:  # BytesIO object
            # Read first few bytes to detect file signature
            file_bytes = file_data.getvalue()
            file_data.seek(0)  # Reset position
            
            # File signature detection
            if len(file_bytes) >= 4:
                # PDF signature: %PDF
                if file_bytes[:4] == b'%PDF':
                    return 'application/pdf'
                
                # PNG signature: \x89PNG\r\n\x1a\n
                if file_bytes[:8] == b'\x89PNG\r\n\x1a\n':
                    return 'image/png'
                
                # JPEG signature: \xff\xd8\xff
                if file_bytes[:3] == b'\xff\xd8\xff':
                    return 'image/jpeg'
                
                # MP3 signature: ID3 or \xff\xfb or \xff\xf3 or \xff\xf2
                if (file_bytes[:3] == b'ID3' or 
                    file_bytes[:2] in [b'\xff\xfb', b'\xff\xf3', b'\xff\xf2'] or
                    (len(file_bytes) >= 10 and file_bytes[:10] == b'ID3')):
                    return 'audio/mpeg'
                
                # WAV signature: RIFF....WAVE
                if file_bytes[:4] == b'RIFF' and file_bytes[8:12] == b'WAVE':
                    return 'audio/wav'
                
                # Video format detection (more comprehensive)
                # MP4 signature: ftyp
                if file_bytes[:4] == b'ftyp':
                    # Check for specific MP4 types
                    if len(file_bytes) >= 8:
                        ftyp_data = file_bytes[4:8]
                        if b'mp4' in ftyp_data or b'M4V' in ftyp_data or b'isom' in ftyp_data:
                            return 'video/mp4'
                        elif b'M4A' in ftyp_data:
                            return 'audio/mp4'
                        else:
                            # Generic ftyp - could be video or audio
                            return 'video/mp4'  # Default to video
                
                # AVI signature: RIFF....AVI
                if file_bytes[:4] == b'RIFF' and file_bytes[8:12] == b'AVI ':
                    return 'video/avi'
                
                # MOV signature: ftyp (same as MP4 but different container)
                if file_bytes[:4] == b'ftyp':
                    return 'video/quicktime'
                
                # Additional video format detection
                # WebM signature: \x1a\x45\xdf\xa3
                if file_bytes[:4] == b'\x1a\x45\xdf\xa3':
                    return 'video/webm'
                
                # FLV signature: FLV
                if file_bytes[:3] == b'FLV':
                    return 'video/x-flv'
                
                # MKV signature: \x1a\x45\xdf\xa3 (same as WebM)
                if file_bytes[:4] == b'\x1a\x45\xdf\xa3':
                    return 'video/x-matroska'
                
                # 3GP signature: ftyp
                if file_bytes[:4] == b'ftyp' and len(file_bytes) >= 8:
                    if b'3gp' in file_bytes[4:8]:
                        return 'video/3gpp'
                
                # DOCX signature: PK\x03\x04 (ZIP format)
                if file_bytes[:4] == b'PK\x03\x04':
                    return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                
                # PPTX signature: PK\x03\x04 (ZIP format)
                if file_bytes[:4] == b'PK\x03\x04':
                    return 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                
                # Additional audio format detection
                # OGG signature: OggS
                if file_bytes[:4] == b'OggS':
                    return 'audio/ogg'
                
                # FLAC signature: fLaC
                if file_bytes[:4] == b'fLaC':
                    return 'audio/flac'
                
                # AAC signature: ADIF or ADTS
                if file_bytes[:4] == b'ADIF' or file_bytes[:2] == b'\xff\xf1':
                    return 'audio/aac'
            
            # Try to detect text files
            try:
                text_content = file_bytes.decode('utf-8', errors='ignore')
                if text_content.isprintable() and len(text_content) > 0:
                    return 'text/plain'
            except:
                pass
                
        return None
    except Exception as e:
        update_progress(f"Error detecting file type: {str(e)}")
        return None

def get_file_type(file_data, file_meta=None, file_name=None):
    """Get file type using MIME detection first, then extension fallback"""
    
    # 1. Try MIME detection first (most reliable)
    detected_mime = detect_file_type(file_data, file_meta)
    # update_progress(f"File type detection for '{file_name}': MIME={detected_mime}")
    
    # 2. If MIME detection works, use it
    if detected_mime:
        mime_to_extension = {
            'text/plain': '.txt',
            'application/pdf': '.pdf',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
            'application/msword': '.doc',
            'audio/mpeg': '.mp3',
            'audio/wav': '.wav',
            'audio/ogg': '.ogg',
            'audio/flac': '.flac',
            'audio/aac': '.aac',
            'audio/mp4': '.m4a',
            'video/mp4': '.mp4',
            'video/avi': '.avi',
            'video/quicktime': '.mov',
            'video/webm': '.webm',
            'video/x-flv': '.flv',
            'video/x-matroska': '.mkv',
            'video/3gpp': '.3gp',
            'image/jpeg': '.jpg',
            'image/png': '.png',
            'application/vnd.ms-powerpoint': '.ppt',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx'
        }
        if detected_mime in mime_to_extension:
            # update_progress(f"Detected MIME type: {detected_mime} for file '{file_name}'")
            return mime_to_extension[detected_mime], detected_mime
    
    # 3. Fallback to extension-based detection
    if file_name:
        file_extension = os.path.splitext(file_name)[1].lower()
        if file_extension in ['.txt', '.pdf', '.docx', '.doc', '.mp3', '.mp4', '.wav', '.avi', '.mov', '.jpg', '.jpeg', '.png', '.ppt', '.pptx', '.ogg', '.flac', '.aac', '.m4a', '.webm', '.flv', '.mkv', '.3gp']:
            return file_extension, None
    
    # 4. For files without extension, try to guess based on MIME type or file content
    if detected_mime:
        # If we have a MIME type but no extension mapping, try to infer
        if 'video' in detected_mime:
            #update_progress(f"Inferring video file type from MIME '{detected_mime}' for '{file_name}'")
            return '.mp4', detected_mime  # Default to mp4 for video files
        elif 'audio' in detected_mime:
            #update_progress(f"Inferring audio file type from MIME '{detected_mime}' for '{file_name}'")
            return '.mp3', detected_mime  # Default to mp3 for audio files
        elif 'image' in detected_mime:
            #update_progress(f"Inferring image file type from MIME '{detected_mime}' for '{file_name}'")
            return '.jpg', detected_mime  # Default to jpg for image files
    
    update_progress(f"Could not determine file type for '{file_name}' (MIME: {detected_mime})")
    return None, detected_mime

def import_media_libs():
    """Lazy import for ffmpeg, whisper, and resampy to improve startup time. Ensures ffmpeg is imported after FFMPEG_BINARY is set."""
    try:
        import importlib
        ffmpeg = importlib.import_module("ffmpeg")
        whisper = importlib.import_module("whisper")
        resampy = importlib.import_module("resampy")
        return ffmpeg, whisper, resampy
    except ImportError as e:
        update_progress(f"Error importing media libraries: {str(e)}")
        return None, None, None
    
thread_pool = ThreadPoolExecutor() # Initialize thread pool for async operations

def init_pytorch():
    try:
        # Only configure PyTorch, do not create or assign torch.classes
        torch.set_default_dtype(torch.float32)
        device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
        torch.set_default_device(device)        
    except Exception as e:
        print(f"PyTorch initialization warning: {str(e)}")

def setup_asyncio():
    try:
        # Apply nest_asyncio first
        nest_asyncio.apply()
        
        # Create new event loop and set as default
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        
        return loop
    except Exception as e:
        print(f"Asyncio setup warning: {str(e)}")
        return None

# Ensure initialization happens at startup
init_pytorch()
loop = setup_asyncio()

# Decorator for running async functions
def run_async(func):
    @functools.wraps(func)
    def wrapper(*args, **kwargs):
        try:
            loop = asyncio.get_event_loop()
        except RuntimeError:
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
        
        return loop.run_until_complete(func(*args, **kwargs))
    return wrapper

# Set up Google API credentials
def setup_google_api():
    # Load environment variables
    load_dotenv(dotenv_path=Path('.env'))  
    try:
        api_key = None
        try:
            api_key = st.secrets["GOOGLE_API_KEY"]
        except Exception:
            api_key = os.getenv("GOOGLE_API_KEY")
        if not api_key:
            st.Error("Google API key not found in secrets or environment variables")
            return False
        Ggenai.configure(api_key=api_key)
        #st.session_state.llm = Ggenai.GenerativeModel(model_name="gemini-2.0-flash-thinking-exp-01-21")
        st.session_state.llm = Ggenai.GenerativeModel(model_name="gemini-2.5-flash-preview-05-20")
    except Exception as e:
        update_progress(f"Error configuring Google Gemini API: {e}")
        st.session_state.llm = None
        return False 
    # Google Drive API setup
    credentials = None
    try:
        credentials = service_account.Credentials.from_service_account_info(
            st.secrets["GCP_SERVICE_ACCOUNT"], scopes=['https://www.googleapis.com/auth/drive']
        )
    except Exception: # Fall back to local service account file
        SERVICE_ACCOUNT_FILE = 'D:/PROJECT/Cupertino/CSTU_Startup_Google_Service_Account.json'
        if os.path.exists(SERVICE_ACCOUNT_FILE):
            try:
                credentials = service_account.Credentials.from_service_account_file(
                    SERVICE_ACCOUNT_FILE, scopes=['https://www.googleapis.com/auth/drive']
                )
            except Exception as e:
                st.error(f"Error loading service account file: {e}")
                return False
    if credentials:
        try:
            st.session_state.google_drive_service = build('drive', 'v3', credentials=credentials)            
            # Test the connection
            test_query = st.session_state.google_drive_service.files().list(
                pageSize=1,
                fields="files(id, name)",
                supportsAllDrives=True,
                includeItemsFromAllDrives=True
            ).execute()        
            if test_query:
                update_progress("Successfully connected to Google Drive")
                return True
            else:
                st.warning("Connected to Google Drive but no access to files")
                return False
        except Exception as e:
            update_progress(f"Error connecting to Google Drive: {str(e)}")
            st.session_state.google_drive_service = None
            return False
    else:
        st.warning("No valid Google Drive credentials found. Drive features will be disabled.")
        st.session_state.google_drive_service = None
        return False

def get_drive_folder_name_by_id(folder_id):
    try:
        folder = st.session_state.google_drive_service.files().get(
            fileId=folder_id,
            fields="name"
        ).execute()
        return folder.get('name', '')
    except Exception as e:
        update_progress(f"Error fetching folder name for id {folder_id}: {str(e)}")
        return ''
    
# Initialize all session state variables
if 'llm' not in st.session_state:
    st.session_state.llm = None
if 'google_drive_service' not in st.session_state:
    st.session_state.google_drive_service = None
if 'save_action' not in st.session_state:
    st.session_state.save_action = None
# Ensure Google API is set up before using get_drive_folder_name_by_id
if not st.session_state.llm or not st.session_state.google_drive_service:
    setup_google_api()    
if 'story_content' not in st.session_state:
    st.session_state.story_content = ""
if 'story_file_stem' not in st.session_state:
    st.session_state.story_file_stem = ""
if 'source_story_root_folder' not in st.session_state:
    st.session_state.source_story_root_folder = {'id': '1fddGiZ4m6m5Gd6_cW9OPr-n01DsPOHEg', 'name': 'Cupertino Stories'}
if 'drive_folder_history' not in st.session_state:
    st.session_state.drive_folder_history = [{'id': 'root', 'name': 'GOOGLE DRIVE'}, {'id': st.session_state.source_story_root_folder['id'], 'name': st.session_state.source_story_root_folder['name']}]
if 'current_folder_files' not in st.session_state:
    st.session_state.current_folder_files = None    
if 'generated_story_root_folder' not in st.session_state:
    folder_id = '1Ld99r4U7x--wdnNgfoPNC67Z6yaqN_A2' # Folder in Google Drive: AI App-Generated Stories
    folder_name = get_drive_folder_name_by_id(folder_id)
    st.session_state.generated_story_root_folder = {'id': folder_id, 'name': folder_name}   
    # st.session_state.story_root_folder = {'id': '1Ld99r4U7x--wdnNgfoPNC67Z6yaqN_A2', 'name': 'MyDrive/Cupertino Stories/AI App that Generate Story from Multiple Files'}
    # st.session_state.story_root_folder = {'id': '1J1QRN0CFaaNpbhTheGyU5h_vYFK-5SpM', 'name': 'MyDrive/PROJECTS'}
# if 'source_folder' not in st.session_state:
    # st.session_state.source_folder = None
if 'progress_history' not in st.session_state:
    st.session_state.progress_history = []
if 'progress_value' not in st.session_state:
    st.session_state.progress_value = 0
if "story_revision_instruction" not in st.session_state:
    st.session_state["story_revision_instruction"] = ""
temp = """
# Create local directories for running locally
DATASOURCE = os.path.join(os.getcwd(), "DATASOURCE")
TEMP = os.path.join(os.getcwd(), "TEMP") # Use for old app
TEMP_GOOGLE_DRIVE = os.path.join(os.getcwd(), "TEMP_GOOGLE_DRIVE")
TEMP_LOCAL_DRIVE = os.path.join(os.getcwd(), "TEMP_LOCAL_DRIVE")
STORIES = os.path.join(os.getcwd(), "STORIES")

os.makedirs(DATASOURCE, exist_ok=True)
os.makedirs(TEMP, exist_ok=True)
os.makedirs(STORIES, exist_ok=True)
os.makedirs(TEMP_GOOGLE_DRIVE, exist_ok=True)
os.makedirs(TEMP_LOCAL_DRIVE, exist_ok=True)
"""
# Configure page
st.set_page_config(
    page_title="CSTU Startup Center - Generate Cupertino Stories",
    page_icon="✍️",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Initialize message queue for thread-safe updates
message_queue = queue.Queue()

def standardize_path(file_path): #Convert any file path to a standardized format with forward slashes
    return os.path.normpath(file_path).replace('\\', '/')

def update_progress_bar(value):
    if 'progress_value' in st.session_state:
        st.session_state.progress_value = value

def ensure_ffmpeg_binary():
    """Ensure FFmpeg binary is available. If not in PATH, download to temp dir (for Streamlit Cloud). Returns ffmpeg binary path or None."""
    import shutil
    import platform
    import stat
    import urllib.request
    import zipfile
    import tarfile
    
    # Check if ffmpeg is in PATH
    ffmpeg_path = shutil.which("ffmpeg")
    if ffmpeg_path:
        return ffmpeg_path
    
    # If not, download static binary
    system = platform.system().lower()
    temp_dir = os.path.join(tempfile.gettempdir(), "ffmpeg-bin")
    os.makedirs(temp_dir, exist_ok=True)
    ffmpeg_bin = os.path.join(temp_dir, "ffmpeg")
    
    if system == "windows":
        ffmpeg_bin += ".exe"
        url = "https://www.gyan.dev/ffmpeg/builds/ffmpeg-release-essentials.zip"
        zip_path = os.path.join(temp_dir, "ffmpeg.zip")
        
        if not os.path.exists(ffmpeg_bin):
            try:
                # update_progress(f"Downloading FFmpeg for Windows from {url}")
                urllib.request.urlretrieve(url, zip_path)                
                # update_progress("Extracting FFmpeg from zip file...")
                with zipfile.ZipFile(zip_path, 'r') as zip_ref:
                    # Find the ffmpeg.exe file in the zip
                    ffmpeg_files = [f for f in zip_ref.namelist() if f.endswith("ffmpeg.exe")]
                    if ffmpeg_files:
                        # Extract the first ffmpeg.exe found
                        zip_ref.extract(ffmpeg_files[0], temp_dir)
                        # The extracted file might be in a subdirectory
                        extracted_path = os.path.join(temp_dir, ffmpeg_files[0])
                        if os.path.exists(extracted_path):
                            shutil.move(extracted_path, ffmpeg_bin)
                            # update_progress(f"FFmpeg extracted to: {ffmpeg_bin}")
                        # else: update_progress(f"Failed to find extracted FFmpeg at: {extracted_path}")
                    # else: update_progress("No ffmpeg.exe found in downloaded zip file")                        
                os.remove(zip_path)
            except Exception as e:
                # update_progress(f"Error downloading/extracting FFmpeg: {str(e)}")
                return None
    else:
        # Assume Linux x86_64 (for Streamlit Cloud)
        url = "https://johnvansickle.com/ffmpeg/releases/ffmpeg-release-amd64-static.tar.xz"
        tar_path = os.path.join(temp_dir, "ffmpeg.tar.xz")
        
        if not os.path.exists(ffmpeg_bin):
            try:
                # update_progress(f"Downloading FFmpeg for Linux from {url}")
                urllib.request.urlretrieve(url, tar_path)                
                #update_progress("Extracting FFmpeg from tar file...")
                with tarfile.open(tar_path, 'r:xz') as tar_ref:
                    for member in tar_ref.getmembers():
                        if member.isfile() and member.name.endswith("/ffmpeg"):
                            tar_ref.extract(member, temp_dir)
                            src = os.path.join(temp_dir, member.name)
                            shutil.move(src, ffmpeg_bin)
                            os.chmod(ffmpeg_bin, os.stat(ffmpeg_bin).st_mode | stat.S_IEXEC)
                            # update_progress(f"FFmpeg extracted to: {ffmpeg_bin}")
                            break
                os.remove(tar_path)
            except Exception as e:
                #update_progress(f"Error downloading/extracting FFmpeg: {str(e)}")
                return None
    
    if os.path.exists(ffmpeg_bin):
        # Test if the binary is executable
        try:
            if system == "windows":
                # On Windows, just check if file exists and has .exe extension
                if ffmpeg_bin.endswith('.exe'):
                    # update_progress(f"FFmpeg binary ready: {ffmpeg_bin}")
                    return ffmpeg_bin
            else:
                # On Linux, check if executable
                if os.access(ffmpeg_bin, os.X_OK):
                    #update_progress(f"FFmpeg binary ready: {ffmpeg_bin}")
                    return ffmpeg_bin
        except Exception as e:
            update_progress(f"Error testing FFmpeg binary: {str(e)}")    
    # update_progress("FFmpeg binary not available")
    return None

def convert_video_to_audio(video_file):
    try:
        # Verify the input file exists
        if not os.path.exists(video_file):
            update_progress(f"Video file does not exist: {video_file}")
            return None
        
        file_size = os.path.getsize(video_file)
        if file_size == 0:
            update_progress(f"Video file is empty: {video_file}")
            return None
            
        #update_progress(f"Processing video file: {video_file} (size: {file_size} bytes)")
        
        ffmpeg_bin = ensure_ffmpeg_binary()
        if ffmpeg_bin:
            # Set FFMPEG_BINARY environment variable for ffmpeg-python
            os.environ["FFMPEG_BINARY"] = ffmpeg_bin
            # Also add to PATH for compatibility
            ffmpeg_dir = os.path.dirname(ffmpeg_bin)
            if ffmpeg_dir not in os.environ["PATH"]:
                os.environ["PATH"] = ffmpeg_dir + os.path.sep + os.environ["PATH"]
            #update_progress(f"Using FFmpeg binary at: {ffmpeg_bin}")
        else:
            update_progress("FFmpeg binary not found, trying system FFmpeg")
            
        ffmpeg, _, _ = import_media_libs() # Import after env var set
        if ffmpeg is None:
            update_progress("Failed to import ffmpeg library")
            return None
            
        #update_progress(f"Starting video to audio conversion...")
        
        # Use explicit ffmpeg binary path if available
        if ffmpeg_bin:
            try:
                # Try using the binary directly with subprocess as fallback                                
                # Create temporary output file
                with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as temp_output:
                    temp_output_path = temp_output.name
                
                #update_progress(f"Using subprocess approach with binary: {ffmpeg_bin}")
                #update_progress(f"Input file: {video_file}")
                #update_progress(f"Output file: {temp_output_path}")
                
                # Run ffmpeg command directly
                cmd = [
                    ffmpeg_bin,
                    '-i', video_file,
                    '-f', 'wav',
                    '-acodec', 'pcm_s16le',
                    '-ac', '1',
                    '-ar', '16000',
                    '-y',  # Overwrite output file
                    temp_output_path
                ]
                
                #update_progress(f"Running command: {' '.join(cmd)}")
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
                
                #update_progress(f"Subprocess return code: {result.returncode}")
                #if result.stderr:
                #    update_progress(f"Subprocess stderr: {result.stderr}")
                #if result.stdout:
                #    update_progress(f"Subprocess stdout: {result.stdout}")
                
                if result.returncode == 0 and os.path.exists(temp_output_path):
                    # Read the output file
                    with open(temp_output_path, 'rb') as f:
                        audio_buffer = io.BytesIO(f.read())
                    audio_buffer.seek(0)
                    
                    # Clean up temp file
                    os.unlink(temp_output_path)
                    
                    #update_progress(f"Successfully converted {video_file} to audio stream (size: {len(audio_buffer.getvalue())} bytes)")
                    return audio_buffer
                else:
                    update_progress(f"FFmpeg subprocess failed: {result.stderr}")
                    if os.path.exists(temp_output_path):
                        os.unlink(temp_output_path)
                    
            except Exception as subprocess_error:
                update_progress(f"Subprocess approach failed: {str(subprocess_error)}")
                # Fall back to ffmpeg-python
                pass
        
        # Fallback to ffmpeg-python approach
        try:
            #update_progress("Trying ffmpeg-python approach...")
            audio_stream = ffmpeg.input(video_file).audio
            audio_buffer = io.BytesIO()
            audio_stream = audio_stream.output('pipe:', format='wav', acodec='pcm_s16le', ac=1, ar=16000, loglevel='error')
            out, _ = audio_stream.run(capture_stdout=True, capture_stderr=True)
            audio_buffer.write(out)
            audio_buffer.seek(0)
            #update_progress(f"Successfully converted {video_file} to audio stream (size: {len(audio_buffer.getvalue())} bytes)")
            return audio_buffer
        except Exception as ffmpeg_python_error:
            update_progress(f"FFmpeg-python approach failed: {str(ffmpeg_python_error)}")
            return None
            
    except Exception as e:
        update_progress(f"FFmpeg error: {str(e)}")
        return None

def extract_content(file_data, file_meta=None):
    """Extract content from file data in memory
    Args:
        file_data: Either a file path string or BytesIO object containing file data
        file_meta: Dictionary containing file metadata (name, mimeType, etc.)
    """
    file_content_list = []

    # Get file information
    if isinstance(file_data, str):  # Local file path
        file_name = os.path.basename(file_data)
        file_path = file_data
        is_memory_file = False
    else:  # BytesIO object from Google Drive
        file_name = file_meta.get('name', 'unknown')
        file_path = None
        is_memory_file = True

    # Get file type using hybrid approach
    file_type, detected_mime = get_file_type(file_data, file_meta, file_name)
    
    if not file_type:
        update_progress(f"Unsupported file format for: '{file_name}' (detected MIME: {detected_mime})")
        return file_content_list

    try:
        # Special handling: if file is a Google Doc, treat as docx
        if file_meta and file_meta.get('mimeType', '') == 'application/vnd.google-apps.document':
            try:
                doc = docx.Document(file_data)
                for paragraph in doc.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        file_content_list.append({"text": text})
                        # update_progress(f"Content extracted successfully from '{file_name}' (Google Doc)")
                doc = None
            except Exception as e:
                update_progress(f"Error processing Google Doc '{file_name}' due to: {str(e)}")
            return file_content_list
        
        # Handle media files (video and audio)
        if file_type in ['.mp4', '.avi', '.mov', '.webm', '.flv', '.mkv', '.3gp', '.mp3', '.wav', '.ogg', '.flac', '.aac', '.m4a']:
            _, whisper, resampy = import_media_libs()  # Get whisper and resampy from lazy import
            if file_type in ['.mp4', '.avi', '.mov', '.webm', '.flv', '.mkv', '.3gp']: # Convert video to audio
                try:
                    if is_memory_file:
                        # For memory file, write to temporary file just for ffmpeg
                        # update_progress(f"Creating temporary file for video processing: {file_name}")
                        
                        # Use a more reliable temporary file creation
                        temp_dir = tempfile.gettempdir()
                        temp_filename = f"video_{hash(file_name)}_{int(time.time())}{file_type}"
                        temp_video = os.path.join(temp_dir, temp_filename)
                        
                        try:
                            with open(temp_video, 'wb') as tmp:
                                file_bytes = file_data.getvalue()
                                tmp.write(file_bytes)
                                tmp.flush()  # Ensure data is written
                                os.fsync(tmp.fileno())  # Force sync to disk
                            
                            #update_progress(f"Temporary file created: {temp_video} (size: {len(file_bytes)} bytes)")
                            
                            # Verify the temporary file exists and has content
                            if os.path.exists(temp_video):
                                actual_size = os.path.getsize(temp_video)
                                # update_progress(f"Temporary file verified: {temp_video} (actual size: {actual_size} bytes)")
                                
                                if actual_size > 0:
                                    #update_progress(f"Converting video to audio: {temp_video}")
                                    audio_source = convert_video_to_audio(temp_video)
                                else:
                                    #update_progress(f"Temporary file is empty: {temp_video}")
                                    audio_source = None
                            else:
                                #update_progress(f"Temporary file does not exist after creation: {temp_video}")
                                audio_source = None
                                
                        except Exception as write_error:
                            update_progress(f"Error writing temporary file: {str(write_error)}")
                            audio_source = None
                        
                        # Clean up temporary file AFTER processing
                        try:
                            if os.path.exists(temp_video):
                                os.unlink(temp_video)
                                # update_progress(f"Cleaned up temporary file: {temp_video}")
                        except Exception as cleanup_error:
                            update_progress(f"Error cleaning up temporary file: {str(cleanup_error)}")
                    else:
                        # update_progress(f"Converting video to audio: {file_path}")
                        audio_source = convert_video_to_audio(file_path)
                        
                    if audio_source is None:
                        update_progress(f"Video to audio conversion failed for '{file_name}'")
                        return file_content_list
                        
                except Exception as e:
                    update_progress(f"Error converting video to audio: {e}\n")
                    return file_content_list
            else: # For audio files 
                if is_memory_file:
                    audio_source = io.BytesIO(file_data.getvalue())
                    audio_source.seek(0) 
                else: # For local files, use the file path directly
                    audio_source = file_path                       
            try: # Process audio content
                audio_data, sample_rate = sf.read(audio_source)                                                    
                if audio_data is None:
                    update_progress(f"Failed to read audio data from '{file_name}'")
                    return file_content_list
                if len(audio_data.shape) > 1:  # If stereo, convert to mono
                    audio_data = audio_data.mean(axis=1)
     
                # Ensure sample rate is 16000 Hz (Whisper expects 16kHz)
                if sample_rate != 16000:
                    try:
                            # Add padding if needed for very short audio
                            if len(audio_data) < sample_rate // 100:  # Less than 10ms
                                pad_length = sample_rate // 100 - len(audio_data)
                                audio_data = np.pad(audio_data, (0, pad_length), mode='constant')
                                
                            audio_data = resampy.resample(audio_data, sample_rate, 16000)
                            sample_rate = 16000
                            
                            if len(audio_data) < 1:
                                # update_progress(f"Resampling resulted in empty audio for {file_name}")
                                return file_content_list
                                
                    except Exception as e:
                            update_progress(f"Resampling failed for '{file_name}' due to: {str(e)}")
                            return file_content_list
                            
                audio_data = audio_data.astype(np.float32)                                         
                device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
                model = whisper.load_model("base").to(device)
                with torch.inference_mode():
                    torch.set_grad_enabled(False)
                    result = model.transcribe(
                        audio_data,
                        fp16=False if device.type == 'cpu' else True
                    ) 
                # update_progress(f"Transcription is completed for '{file_name}'.\n")
                plain_transcript = result["text"].strip() # Plain transcript
                file_content_list.append({"text": plain_transcript}) # Plain transcript
                #update_progress(f"Content extracted successfully from '{file_name}'")
                timestamped_transcript = [] # Create timestamped transcript
                for segment in result["segments"]:
                    start_time = str(datetime.timedelta(seconds=round(segment["start"])))
                    text = segment["text"].strip()
                    timestamped_transcript.append(f"[{start_time}] {text}")    
                #update_progress(f"Successfully processed video/audio file: {file_name}")                                
            except Exception as e:
                update_progress(f"Error processing video/audio file '{file_name}': {str(e)}")
                return file_content_list
            if isinstance(audio_source, io.BytesIO): # Not required
                audio_source.close()

        # Handle images
        elif file_type in ['.jpg', '.jpeg', '.png']:
            try:  # Open image from memory or file
                if is_memory_file:
                    image = Image.open(file_data)
                else:
                    with open(file_path, "rb") as img_file:
                        image = Image.open(img_file)                
                if image.mode in ('RGBA', 'LA', 'P'):
                    rgb_image = image.convert('RGB')
                else:
                    rgb_image = image.copy()                
                buffered = io.BytesIO()
                rgb_image.save(buffered, format="JPEG")
                image_data = buffered.getvalue()                
                mime_type = "image/jpeg"
                if file_type == '.png':
                    mime_type = "image/png"                
                file_content_list.append({"mime_type": mime_type, "data": image_data})  
                #update_progress(f"Content extracted successfully from '{file_name}'")              
                # Clean up
                image.close()
                rgb_image.close()
                buffered.close()
                image_data = None                
            except Exception as e:
                update_progress(f"Error processing image '{file_name}' due to: {str(e)}")
    
        # Handle PDFs
        elif file_type == '.pdf':
            try:
                if is_memory_file: # Load PDF from memory
                    doc = fitz.open(stream=file_data.getvalue(), filetype="pdf")
                else:
                    doc = fitz.open(file_path)                
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    text = page.get_text().strip()
                    if text:
                        file_content_list.append({"text": text})
                        # update_progress(f"Content extracted successfully from '{file_name}'")
                doc.close()                
            except Exception as e:
                update_progress(f"Error processing PDF '{file_name}' due to: {str(e)}")
        
        # Handle Word documents
        elif file_type in ['.doc', '.docx']:
            try:
                if is_memory_file:
                    doc = docx.Document(file_data)
                else:
                    doc = docx.Document(file_path)                
                for paragraph in doc.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        file_content_list.append({"text": text})
                        # update_progress(f"Content extracted successfully from '{file_name}'")
                doc = None                
            except Exception as e:
                update_progress(f"Error processing Word document '{file_name}' due to: {str(e)}")
        
        # Handle text files
        elif file_type == '.txt':
            try:
                if is_memory_file:
                    content = file_data.getvalue().decode('utf-8').strip()
                else:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content = f.read().strip()                
                if content:
                    file_content_list.append({"text": content})
                    # update_progress(f"Content extracted successfully from '{file_name}'")
            except Exception as e:
                update_progress(f"Error processing text file '{file_name}' due to: {str(e)}")

        # Handle PowerPoint files
        elif file_type in ['.ppt', '.pptx']:
            try:
                if is_memory_file:
                    ppt = pptx.Presentation(file_data)
                else:
                    ppt = pptx.Presentation(file_path)
                for slide_num, slide in enumerate(ppt.slides, 1):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text = shape.text.strip()
                            if text:
                                slide_text.append(text)
                    if slide_text:
                        file_content_list.append({"text": f"[Slide {slide_num}] " + "\n".join(slide_text)})
                #update_progress(f"Content extracted successfully from PowerPoint '{file_name}'")
            except Exception as e:
                update_progress(f"Error processing PowerPoint file '{file_name}' due to: {str(e)}")
            return file_content_list        
        else:
            update_progress(f"Unsupported file format for: '{file_name}'")        
            
        # Fallback: If no file type was detected but the file might be audio/video, try to process it
        if not file_content_list and detected_mime and ('audio' in detected_mime or 'video' in detected_mime):
            # update_progress(f"Attempting to process '{file_name}' as media file (MIME: {detected_mime})")
            try:
                _, whisper, resampy = import_media_libs()
                
                # For video files, convert to audio first
                if 'video' in detected_mime:
                    if is_memory_file:
                        # Create temporary file with appropriate extension
                        temp_dir = tempfile.gettempdir()
                        temp_filename = f"video_{hash(file_name)}_{int(time.time())}.mp4"
                        temp_video = os.path.join(temp_dir, temp_filename)
                        
                        try:
                            with open(temp_video, 'wb') as tmp:
                                file_bytes = file_data.getvalue()
                                tmp.write(file_bytes)
                                tmp.flush()
                                os.fsync(tmp.fileno())
                            
                            #update_progress(f"Created temporary video file: {temp_video}")
                            audio_source = convert_video_to_audio(temp_video)
                            
                            # Clean up
                            if os.path.exists(temp_video):
                                os.unlink(temp_video)
                        except Exception as e:
                            update_progress(f"Error creating temporary video file: {str(e)}")
                            audio_source = None
                    else:
                        audio_source = convert_video_to_audio(file_path)
                else:
                    # For audio files, use directly
                    if is_memory_file:
                        audio_source = io.BytesIO(file_data.getvalue())
                        audio_source.seek(0)
                    else:
                        audio_source = file_path
                
                if audio_source is None:
                    update_progress(f"Failed to convert video to audio for '{file_name}'")
                    return file_content_list
                
                audio_data, sample_rate = sf.read(audio_source)
                if audio_data is not None:
                    if len(audio_data.shape) > 1:
                        audio_data = audio_data.mean(axis=1)
                    
                    # Resample if needed
                    if sample_rate != 16000:
                        try:
                            audio_data = resampy.resample(audio_data, sample_rate, 16000)
                            sample_rate = 16000
                        except Exception as e:
                            update_progress(f"Resampling failed: {str(e)}")
                            return file_content_list
                    
                    audio_data = audio_data.astype(np.float32)
                    device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
                    model = whisper.load_model("base").to(device)
                    with torch.inference_mode():
                        torch.set_grad_enabled(False)
                        result = model.transcribe(audio_data, fp16=False if device.type == 'cpu' else True)
                    
                    plain_transcript = result["text"].strip()
                    file_content_list.append({"text": plain_transcript})
                    # update_progress(f"Successfully processed '{file_name}' as media file")
                else:
                    update_progress(f"Failed to read audio data from '{file_name}'")
            except Exception as e:
                update_progress(f"Error processing '{file_name}' as media: {str(e)}")
        
        return file_content_list
    except Exception as e:
        update_progress(f"Error processing '{file_name}' due to: {str(e)}")
        return []
    finally: # Ensure memory cleanup
        gc.collect()

def create_or_get_drive_folder(parent_folder_id, folder_name):
    try:
        query = f"'{parent_folder_id}' in parents and name='{folder_name}' and mimeType='application/vnd.google-apps.folder' and trashed=false"
        list_request = st.session_state.google_drive_service.files().list(
            q=query,
            spaces='drive',
            fields='files(id, name)'
        )
        results = list_request.execute()
        files = results.get('files', [])
        if files:
            return files[0]['id']
        folder_metadata = {
            'name': folder_name,
            'mimeType': 'application/vnd.google-apps.folder',
            'parents': [parent_folder_id]
        }
        create_request = st.session_state.google_drive_service.files().create(
            body=folder_metadata,
            fields='id'
        )
        folder = create_request.execute()
        return folder.get('id')
    except Exception as e:
        update_progress(f"Error creating/getting folder in Google Drive: {str(e)}")
        return parent_folder_id

def check_file_exists_in_drive(folder_id, filename):
    """Return file ID if file exists in folder, else None."""
    try:
        query = f"'{folder_id}' in parents and name='{filename}' and trashed=false"
        list_request = st.session_state.google_drive_service.files().list(
            q=query,
            fields='files(id, name)',
            spaces='drive'
        )
        results = list_request.execute()
        files = results.get('files', [])
        if files:
            return files[0]['id']  # Return file ID if found
        return None
    except Exception as e:
        update_progress(f"Error checking file existence: {str(e)}")
        return None

def markdown_to_docx(md_text, doc):
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Arial'
    font.size = Pt(10)
    # For compatibility with some Word versions
    style.element.rPr.rFonts.set(qn('w:eastAsia'), 'Arial')
    lines = md_text.splitlines()

    for line in lines:
        line = line.strip()
        if line.startswith('# '):
            h = doc.add_heading(line[2:].strip(), level=1)
            h.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
            run = h.runs[0]
            run.font.name = 'Arial'
            run.font.size = Pt(15)
            run.bold = True
            h.style.font.name = 'Arial'
            h.style.font.size = Pt(15)
        elif line.startswith('## '):
            h = doc.add_heading(line[3:].strip(), level=2)
            h.alignment = WD_PARAGRAPH_ALIGNMENT.JUSTIFY
            run = h.runs[0]
            run.font.name = 'Arial'
            run.font.size = Pt(12)
            run.bold = True
            h.style.font.name = 'Arial'
            h.style.font.size = Pt(12)
        elif line.startswith('### '):
            h = doc.add_heading(line[4:].strip(), level=3)
            h.alignment = WD_PARAGRAPH_ALIGNMENT.JUSTIFY
            run = h.runs[0]
            run.font.name = 'Arial'
            run.font.size = Pt(11)
            run.bold = True
            h.style.font.name = 'Arial'
            h.style.font.size = Pt(11)
        elif line:
            para = doc.add_paragraph(line)
            para.alignment = WD_PARAGRAPH_ALIGNMENT.JUSTIFY
            for run in para.runs:
                run.font.name = 'Arial'
                run.font.size = Pt(10)

def save_story_to_google_drive():
    if st.session_state.google_drive_service:
        try: # Use drive_folder_history to get story_folder_name
            current_folder = st.session_state.drive_folder_history[-1]
            current_folder_id = current_folder['id']
            story_folder_name = current_folder['name']
            story_folder_id = create_or_get_drive_folder(st.session_state.generated_story_root_folder['id'], story_folder_name)
            if st.session_state.save_action is None:
                docx_exists = check_file_exists_in_drive(story_folder_id, os.path.basename(st.session_state.story_file_stem) + '.docx')
                html_exists = check_file_exists_in_drive(story_folder_id, os.path.basename(st.session_state.story_file_stem) + '.html')
                if docx_exists or html_exists: 
                    st.session_state.overwrite_confirm_needed = True
                    st.session_state.existing_txt_file_id = docx_exists
                    st.session_state.existing_html_file_id = html_exists
                    return
                else:
                    st.session_state.save_action = 'save'

            if st.session_state.save_action in ['save', 'save_with_timestamp']:                 
                if st.session_state.save_action == 'save_with_timestamp':
                    unique_suffix = datetime.datetime.now().strftime("_%Y%m%d_%H%M%S")
                else: 
                    unique_suffix = ''
                st.session_state.story_file_stem = st.session_state.story_file_stem + unique_suffix
                # Save as .docx instead of .txt
                # md_text = st.session_state.story_content
                # html_body = markdown2.markdown(md_text)
                # Fallback: ensure at least <p> or <h> tags for html2docx
                # if '<p>' not in html_body and '<h' not in html_body:
                #    html_body = ''.join(f'<p>{line}</p>' for line in md_text.splitlines() if line.strip())
                doc = docx.Document()
                markdown_to_docx(st.session_state.story_content, doc)
                # html2docx(html_body, doc)         
                docx_buffer = io.BytesIO()
                doc.save(docx_buffer)
                docx_buffer.seek(0)
                docx_metadata = {
                    'name': f"{st.session_state.story_file_stem}.docx",
                    'parents': [story_folder_id],
                    'mimeType': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                }
                docx_media = MediaIoBaseUpload(
                    docx_buffer,
                    mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    resumable=True
                )
                st.session_state.google_drive_service.files().create(
                    body=docx_metadata,
                    media_body=docx_media,
                    fields='id'
                ).execute()
                html_metadata = {
                    'name': f"{st.session_state.story_file_stem}.html",
                    'parents': [story_folder_id],
                    'mimeType': 'text/html'
                }
                html_media = MediaIoBaseUpload(
                    io.BytesIO(story_html_format().encode('utf-8')),
                    mimetype='text/html',
                    resumable=True
                )
                st.session_state.google_drive_service.files().create(
                    body=html_metadata,
                    media_body=html_media,
                    fields='id'
                ).execute()

            if st.session_state.save_action == 'overwrite':
                docx_file_id = st.session_state.get("existing_txt_file_id")  # Now used for .docx
                html_file_id = st.session_state.get("existing_html_file_id")
                # Overwrite .docx
                if docx_file_id:
                    doc = docx.Document()
                    markdown_to_docx(st.session_state.story_content, doc)
                    docx_buffer = io.BytesIO()
                    doc.save(docx_buffer)
                    docx_buffer.seek(0)
                    docx_media = MediaIoBaseUpload(
                        docx_buffer,
                        mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                        resumable=True
                    )
                    st.session_state.google_drive_service.files().update(
                        fileId=docx_file_id,
                        media_body=docx_media
                    ).execute()
                # Overwrite .html
                if html_file_id:
                    html_media = MediaIoBaseUpload(
                        io.BytesIO(story_html_format().encode('utf-8')),
                        mimetype='text/html',
                        resumable=True
                    )
                    st.session_state.google_drive_service.files().update(
                        fileId=html_file_id,
                        media_body=html_media
                    ).execute()            
            update_progress(f"Story saved as '{st.session_state.story_file_stem}.docx' and '{st.session_state.story_file_stem}.html' in '{st.session_state.generated_story_root_folder['name']}'/'{story_folder_name}' folder.")
            st.session_state.current_folder_files = None  # Reset current folder files to force reload file list      
        except Exception as e:
            update_progress(f"Error saving to Google Drive: {str(e)}")
            st.error(f"Error saving to Google Drive: {str(e)}")            
            raise  # Re-raise the exception to be handled by the caller
    st.session_state.save_action = None
    return
def story_html_format():
    html_body = markdown.markdown(st.session_state.story_content)
    story_content_html = f"""
        <html>
        <head>
            <style>
                body {{font-family: Arial, sans-serif; font-size: 10; margin: 10px; text-align: justify; padding: 0;}}
                h1 {{font-family: Arial, sans-serif; font-size: 15pt; text-align: center; font-weight: bold;}}
                h2 {{font-family: Arial, sans-serif; font-size: 12pt; text-align: justify; font-weight: bold;}}
                h3 {{font-family: Arial, sans-serif; font-size: 11pt; text-align: justify; font-weight: bold;}}
            </style>
        </head>
        <body>
            {html_body}
        </body>
        </html>"""
    return story_content_html

async def generate_story(file_content_list):
    """Generate story from content and save to both local and Google Drive"""
    if not st.session_state.get('llm'):
        update_progress("LLM not initialized. Cannot generate story.")
        return False        
    selected_file_names = []
    for file in st.session_state.selected_files:
        selected_file_names.append(f"'{file['name']}'")
    selected_file_names = ', '.join(selected_file_names)      
    prompt = ("Generate story content without mentioning this prompt in English then Chinese but never interleaved. " 
              "The structure must follow: title in heading 1, Authors (heading 2): include author, interviewer or organization name and contact info if available (body text), "
              "Category (heading 2): If the generated story is related to Cupertino, select one or more categories from: Education, Community, Environment, Government, Business, or Safety (body text). Else set category to 'Unrelated to Cupertino' "
              "Summary (heading 2): include what story is about (body text), Full Story (heading 2): complete story formatted for easy reading. The data source is below:")
    try:       
        response = await asyncio.to_thread(st.session_state.llm.generate_content, [prompt] + file_content_list)
        st.session_state.story_content = response.text #story_content_markdown               
        update_progress(f"Story generated from {selected_file_names}")                  
        save_story_to_google_drive()
    except Exception as e:
        update_progress(f"Error generating story from {selected_file_names} due to: {str(e)}")
    finally:     
        gc.collect() # Clean up
        return

def revise_story():
    if not st.session_state.get('llm'):
        update_progress("LLM not initialized. Cannot generate story.")
        return False             
    # prompt = "Generate content without mentioning this prompt in English then Chinese but never interleaved. The structure must follow: title in heading 1, Authors (heading 2): include author, interviewer or organization name and contact info if available (body text), Category (heading 2): selected one or more from Education/Community/Environment/Government/Business/Public Safety (body text), Summary (heading 2): include what story is about (body text), Full Story (heading 2): complete story formatted for easy reading. The data source is below:"
    try:       
        response = st.session_state.llm.generate_content(f"Here is story: '{st.session_state.story_content}'. Try to keep the structure and revise story content with instruction below: '{st.session_state.story_revision_instruction}'")
        st.session_state.story_content = response.text #story_content_markdown               
    except Exception as e:
        update_progress(f"Error revising story from current content due to: {str(e)}")
    return

async def process_local_files(uploaded_files):
    if not uploaded_files:
        st.warning("No files selected.")
        return
hidden ="""
    for uploaded_file in uploaded_files:
        source_file_name = uploaded_file.name
        update_progress(f"Processing {source_file_name}...")        
        st.session_state.source_files.append(os.path.join(TEMP_LOCAL_DRIVE, source_file_name))
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(source_file_name)[1]) as tmp:
            tmp.write(uploaded_file.getvalue())
            temp_path = standardize_path(tmp.name)
            st.session_state.source_files.append((temp_path, source_file_name))
            
    if process_mode == 'per_file':
        for temp_path, orig_path in st.session_state.source_files:
            story_file_base_path = standardize_path(os.path.join(STORIES, os.path.splitext(os.path.basename(orig_path))[0]))
            update_progress(f"Extracting content from {orig_path}...")
            file_content_list = extract_content(temp_path)
            if file_content_list:
                await generate_story(file_content_list, orig_path, story_file_base_path)
    else:
        file_content_list = []
        for temp_path, orig_path in st.session_state.source_files:
            update_progress(f"Extracting content from {orig_path}...")
            content = extract_content(temp_path)
            if content:
                file_content_list.extend(content)
        
        if file_content_list:
            folder_name = "Combined_Story_" + datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            story_file_base_path = os.path.join(STORIES, folder_name)
            source_path = "Multiple Files from Local Drive"
            await generate_story(file_content_list, source_path, story_file_base_path)

    # Clean up temp files and metadata
    for temp_path, _ in st.session_state.source_files:
        if os.path.exists(temp_path):
            os.unlink(temp_path)
    st.session_state.source_files = []
"""

def download_google_drive_file(file_id): #"""Download a file from Google Drive and return its content as a BytesIO object"""
    try:
        # Get file metadata to check mimeType
        file_meta = st.session_state.google_drive_service.files().get(fileId=file_id, fields="name, mimeType").execute()
        file_name = file_meta.get("name", file_id)
        mime_type = file_meta.get("mimeType", "")
        # If Google Docs file, use export_media
        if mime_type == "application/vnd.google-apps.document":
            # Export as docx (preferred for further processing)
            export_mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            try:
                request = st.session_state.google_drive_service.files().export_media(fileId=file_id, mimeType=export_mime)
                file_content = io.BytesIO()
                downloader = MediaIoBaseDownload(file_content, request)
                done = False
                while not done:
                    _, done = downloader.next_chunk()
                file_content.seek(0)
                # Patch file_meta to reflect .docx for downstream processing
                file_meta["name"] = file_name + ".docx" if not file_name.lower().endswith(".docx") else file_name
                return file_content
            except Exception as e:
                update_progress(f"Error exporting Google Doc {file_name} as docx: {str(e)}")
                return None
        else:
            # For all other files, use get_media as before
            request = st.session_state.google_drive_service.files().get_media(fileId=file_id)
            file_content = io.BytesIO()
            downloader = MediaIoBaseDownload(file_content, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            file_content.seek(0)
            return file_content
    except Exception as e:
        try:
            # Try to get file name for error reporting
            if 'file_name' not in locals():
                file_meta = st.session_state.google_drive_service.files().get(fileId=file_id, fields="name").execute()
                file_name = file_meta.get("name", file_id)
        except Exception:
            file_name = file_id
        update_progress(f"Error downloading file {file_name} due to: {str(e)}")
        return None

def list_google_drive_files(parent_id='root'):
    """List files and folders in Google Drive with retry logic for SSL errors"""
    if not st.session_state.get('google_drive_service'):
        st.error("Google Drive service is not initialized")
        return []

    max_retries = 5
    backoff = 2
    for attempt in range(max_retries):
        try:
            # Simple query for all files and folders
            query = f"'{parent_id}' in parents and trashed = false"
            if parent_id == 'root':
                query = "('root' in parents or sharedWithMe = true) and trashed = false"
            results = st.session_state.google_drive_service.files().list(
                q=query,
                fields="files(id, name, mimeType, size)",
                pageSize=1000,
                spaces='drive'
            ).execute()
            files = results.get('files', [])
            if files:
                return sorted(files, key=lambda x: (x['mimeType'] != 'application/vnd.google-apps.folder', x['name'].lower()))
            else:
                # Try to get folder name from drive_folder_history
                folder_name = parent_id
                for folder in reversed(st.session_state.drive_folder_history):
                    if folder['id'] == parent_id:
                        folder_name = folder['name']
                        break
                update_progress(f"No files found in folder '{folder_name}'")
                return []
        except (ssl.SSLError, socket.error, HttpError) as e:
            error_str = str(e)
            # Only retry on SSL/connection errors
            if any(msg in error_str for msg in ["SSL", "_ssl.c", "EOF occurred", "record layer failure", "Connection reset", "Connection aborted", "timed out"]):
                update_progress(f"SSL/network error on attempt {attempt+1}: {error_str}. Retrying in {backoff}s...")
                time.sleep(backoff)
                backoff *= 2
                continue
            else:
                error_msg = f"Error listing files: {error_str}"
                st.error(error_msg)
                update_progress(error_msg)
                return []
        except Exception as e:
            error_msg = f"Error listing files: {str(e)}"
            st.error(error_msg)
            update_progress(error_msg)
            return []
    error_msg = f"Google Drive connection failed after {max_retries} attempts. Please check your network or try again later."
    st.error(error_msg)
    update_progress(error_msg)
    return []
def google_drive_browser():
    """Display Google Drive browser in sidebar"""
    # Always reset folder_nav_in_progress at the start so navigation is never stuck
    st.session_state.folder_nav_in_progress = False
    try:
        # Use cached files if available, otherwise fetch new ones
        current_folder = st.session_state.drive_folder_history[-1]
        current_folder_id = current_folder['id']
        if st.session_state.current_folder_files is None:
            st.session_state.current_folder_files = list_google_drive_files(current_folder_id)
        files = st.session_state.current_folder_files
        if 'folder_nav_in_progress' not in st.session_state:
            st.session_state.folder_nav_in_progress = False
        with st.sidebar:
            # Back button if not in root
            if len(st.session_state.drive_folder_history) > 1:                
                if st.button("⬅️**Back to Previous Path**", disabled=st.session_state.folder_nav_in_progress):
                    st.session_state.folder_nav_in_progress = True
                    st.session_state.drive_folder_history.pop()
                    st.session_state.current_folder_files = None
                    time.sleep(0.5)  # Debounce rapid clicks
                    st.session_state.folder_nav_in_progress = False  # Ensure re-enable before rerun
                    st.rerun()                  
            # Label and buttons for parent folders (excluding current folder)
            if len(st.session_state.drive_folder_history) > 1:
                st.markdown("**Parent Folders:**")
                for idx, folder in enumerate(st.session_state.drive_folder_history[:-1]):
                    folder_id = folder['id']
                    folder_name = folder['name']
                    # if st.button(f"📁 {folder_name}", key=f"parent_folder_{folder_id}", help=f"Go to {folder_name}"):
                    if st.checkbox(f"📁 {folder_name}", key=f"parent_folder_{folder_id}", help=f"Go to {folder_name}"):
                        st.session_state.drive_folder_history = st.session_state.drive_folder_history[:idx+1]
                        st.session_state.current_folder_files = None
                        # st.session_state.selected_files = []
                        st.session_state.delete_confirm_needed = False
                        st.session_state.overwrite_confirm_needed = False
                        st.session_state.save_action = None
                        time.sleep(0.5)  # Debounce rapid clicks
                        st.session_state.folder_nav_in_progress = False
                        st.rerun()

            # Current path as plain text
            path_names = [folder['name'] for folder in st.session_state.drive_folder_history]
            st.markdown(f"**Current Folder:** {' > '.join(path_names)}")
            # Label for child folder
            st.markdown("**Child Folders and Files:**")            
            folders = [f for f in files if f['mimeType'] == 'application/vnd.google-apps.folder']
            if folders:
                for folder in sorted(folders, key=lambda x: x['name'].lower()):
                    folder_name = folder['name']
                    folder_id = folder['id']
                    if st.checkbox(f"📁 {folder_name}", key=f"folder_{folder_id}", help="Click to open folder", disabled=st.session_state.folder_nav_in_progress):
                        st.session_state.folder_nav_in_progress = True
                        st.session_state.drive_folder_history.append({'id': folder_id, 'name': folder_name})
                        st.session_state.current_folder_files = None
                        # st.session_state.selected_files = []
                        st.session_state.delete_confirm_needed = False
                        st.session_state.overwrite_confirm_needed = False
                        st.session_state.save_action = None
                        time.sleep(0.5)  # Debounce rapid clicks
                        st.session_state.folder_nav_in_progress = False  # Ensure re-enable before rerun
                        st.rerun()
            if not files:
                st.info("No files found in current folder!")                      
            regular_files = [f for f in files if f['mimeType'] != 'application/vnd.google-apps.folder']
            # Show files with checkboxes            
            all_file_ids = [file['id'] for file in regular_files]
            checked_file_ids = [fid for fid in all_file_ids if st.session_state.get(f"check_{fid}", False)]
            all_selected = len(regular_files) > 0 and len(checked_file_ids) == len(regular_files)
            
            # Track previous state to detect changes
            if 'previous_select_all_state' not in st.session_state:
                st.session_state.previous_select_all_state = all_selected
            
            select_all = st.checkbox("**Select All Files**", 
                value=all_selected,
                key="select_all_files_in_folder", 
                help="Click to select all following files in the current folder")
            
            # Handle select all checkbox changes
            if select_all and not st.session_state.previous_select_all_state:
                # Select all files when checkbox is checked
                for file in regular_files:
                    st.session_state[f"check_{file['id']}"] = True
            elif not select_all and st.session_state.previous_select_all_state:
                # Deselect all files when checkbox is unchecked
                for file in regular_files:
                    st.session_state[f"check_{file['id']}"] = False
            
            # Update previous state
            st.session_state.previous_select_all_state = select_all
            # If html file is selected, update checkbox states before rendering
            # if 'selected_html_file_id' in st.session_state:
            #    for f in regular_files:
            #        st.session_state[f"check_{f['id']}"] = (f['id'] == st.session_state.selected_html_file_id)
            #    del st.session_state['selected_html_file_id'] # Clear html_file_id flag
            st.session_state.selected_files = []
            if regular_files:
                for file in sorted(regular_files, key=lambda x: x['name'].lower()):
                    mime = file.get('mimeType', 'unknown')
                    icon = "📄"  # Default icon
                    if 'image' in mime:
                        icon = "🖼️"
                    elif 'pdf' in mime:
                        icon = "📕"
                    elif 'document' in mime:
                        icon = "📄"
                    elif 'presentation' in mime or 'powerpoint' in mime:
                        icon = "🖥️"
                    elif 'audio' in mime:
                        icon = "🎧"
                    elif 'video' in mime:
                        icon = "🎞️"
                    elif 'html' in mime:
                        icon = "🌐"
                    elif 'spreadsheet' in mime or 'excel' in mime:
                        icon = "📈"                        
                    file_size = int(file.get('size', 0))
                    if file_size >= 1024**3:
                        file_size = f"{round(file_size / 1024**3, 1)} GB"
                    elif file_size >= 1024**2:
                        file_size = f"{round(file_size / 1024**2, 1)} MB"
                    elif file_size >= 1024:
                        file_size = f"{round(file_size / 1024, 1)} KB"
                    else:
                        file_size = f"{file_size} B"

                    col1, col2, col3 = st.columns([2, 50, 6])
                    with col1:
                        file_selected = st.checkbox(" ", key=f"check_{file['id']}")
                    with col2:
                        file_url = f"https://drive.google.com/file/d/{file['id']}/view"
                        st.markdown(f"""
                            <a href="{file_url}" target="_blank" title="Click here to view file in a new tab of browser.\nClick on the left checkbox to select file." style="display:inline-block;padding:0.25em 1em;color:darkblue;border-radius:4px;text-decoration:none;font-size:1em;">
                                {file['name']} [{icon}{file_size}]
                            </a>""", unsafe_allow_html=True)
                    with col3:
                        if file['name'].lower().endswith('.html'):
                            if st.button('🗁', help = 'Click to load and view/edit HTML file in story tabs', key=f"{file['id']}"):
                                file_content = download_google_drive_file(file['id'])
                                if file_content:
                                    html_str = file_content.read().decode('utf-8', errors='replace')
                                    st.session_state.story_content = markdownify.markdownify(html_str, heading_style="ATX")
                                    st.session_state.story_file_stem = file['name'].rsplit('.', 1)[0]
                    # st.success(f"{file['name']} content was loaded to Story Viewer/Editor.")
                    if file_selected:
                        st.session_state.selected_files.append(file)  # Append file metadata             
            if st.session_state.selected_files:   
                if st.button("✍️**GENERATE STORY FROM SELECTED FILES**", help = 'Click to generate a story from selected files', key="generate_from_multi"): 
                    st.session_state.save_action = None
                    st.session_state.overwrite_confirm_needed = False
                    st.session_state.combined_content = []
                    for file_meta in st.session_state.selected_files:
                        file_data = download_google_drive_file(file_meta['id'])
                        if file_data:
                            content = extract_content(file_data, file_meta)
                            if content:
                                st.session_state.combined_content.extend(content)
                    if st.session_state.combined_content:
                        if len(st.session_state.selected_files) == 1:
                            st.session_state.story_file_stem = st.session_state.selected_files[0]['name']
                        else: # Multiple files selected
                            current_folder_id = st.session_state.drive_folder_history[-1]
                            st.session_state.story_file_stem = st.session_state.drive_folder_names.get(current_folder_id, 'GoogleDrive')
                        run_async(generate_story)(st.session_state.combined_content)
                        st.rerun()

                # Overwrite/Skip confirmation UI
                if st.session_state.get('overwrite_confirm_needed', False):
                    current_folder = st.session_state.drive_folder_history[-1]
                    story_folder_name = current_folder['name']
                    st.warning(f"Story files '{st.session_state.story_file_stem}.docx|html' already exist in '{st.session_state.generated_story_root_folder['name']}'/'{story_folder_name}'")
                    options = {
                        "💾Overwrite File": "overwrite",
                        "➕Save with Stamp": "save_with_timestamp",
                        "❌Cancel Saving": "cancel"
                    }
                    selected_option = st.radio(
                        "Choose an action:",
                        list(options.keys()),
                        key="overwrite_radio_option", help = "Overwrite existing story files, Save with timestamp added to existing story files, Cancel saving"
                    )
                    action = options[selected_option]
                    if st.button("✅**Action Confirm**", key="action_confirm"):
                        if action == "overwrite":
                            st.session_state.save_action = 'overwrite'
                            st.session_state.overwrite_confirm_needed = False
                            save_story_to_google_drive()
                            st.rerun()
                        elif action == "save_with_timestamp":
                            st.session_state.save_action = 'save_with_timestamp'
                            st.session_state.overwrite_confirm_needed = False
                            save_story_to_google_drive()
                            st.rerun()
                        elif action == "cancel":
                            st.session_state.save_action = None
                            st.session_state.overwrite_confirm_needed = False
                            update_progress(f"Skipped saving story file '{st.session_state.story_file_stem}'")
                            st.rerun()
                # Delete confirmation UI
                temp = """if "delete_confirm_needed" not in st.session_state:
                    st.session_state.delete_confirm_needed = False                
                if st.button("🗑️**Delete Selected Files**", key="delete_files"):
                    st.session_state.delete_confirm_needed = True
                if st.session_state.delete_confirm_needed:
                    st.warning("Are you sure to delete selected files?")
                    col_yes, col_no = st.columns(2)
                    with col_yes:
                        user_input = st.text_input("Enter authorization code to delete selected files:", type="password")
                        if st.button("Yes, Delete", key="yes_delete"):
                            # today = datetime.now()
                            secret_code = '2025' # f"{today.day:02d}{today.month:02d}"
                            if user_input == secret_code and user_input != "": 
                                try:
                                    for file in list(st.session_state.selected_files): # Use a copy to avoid modifying list during iteration
                                        st.session_state.google_drive_service.files().delete(fileId=file['id']).execute()                                      
                                    st.success("Selected files deleted successfully!")                                   
                                except Exception as e:
                                    st.error(f"Error deleting files: {str(e)}")
                                st.session_state.delete_confirm_needed = False
                                st.session_state.current_folder_files = None  # Reset to refresh the list
                                st.session_state.selected_files = []
                                st.rerun()  # Force rerun to reset checkboxes
                            elif user_input != "":
                                st.error("❌ You are not authorized to delete selected files!")
                    with col_no:
                        if st.button("No, Cancel", key="no_delete"):
                            st.warning("No files deleted!")
                            st.session_state.delete_confirm_needed = False """
    except Exception as e:
        st.error(f"Error in Google Drive browser: {str(e)}")       

@run_async
async def main():
    hide_menu_style = """
        <style>
        #MainMenu {visibility: hidden;}
        </style>
    """
    st.markdown(hide_menu_style, unsafe_allow_html=True)

    # Add error handling for torch path issues
    init_pytorch()

    with st.sidebar:
        st.markdown("""
            <style>
            .sidebar-text {
                font-size: 12px !important;
                color: darkblue !important;
            }       
            div.stButton > button {
                text-align: left !important;
                font-weight: bold !important;
                justify-content: flex-start !important;
                align-items: flex-start !important;
                display: flex !important;
            }                                                   
            .stSelectbox > label,
            .stSelectbox > select {
                color: darkblue;
                font-weight: bold;
            }
            </style>
            """, unsafe_allow_html=True)
        st.markdown("<font color='darkblue'><b><p style='font-size: 16px; text-align: center; '>CSTU Startup Center</p></b></font>", unsafe_allow_html=True)
        st.markdown("<font color='darkblue'><b><p style='font-size: 20px; text-align: center; '>GENERATE CUPERTINO STORIES</p></b></font>", unsafe_allow_html=True)
        st.markdown('<a href="https://docs.google.com/document/d/1dJTJZ58kgweCsS9-lYSRgIRUc160TDnk" target="_blank" title="Click to open the user guide in a new tab" style="display:block; text-align:center; color:darkblue; ">User Guide</a>', unsafe_allow_html=True)
        # st.markdown("---")
  
        #st.header("")
        source = st.radio(label="**Browse Files from:**", options=["Google Drive", "Local Drive"], help="Select source to browse files",)         
        if source == "Local Drive":
            uploaded_files = st.file_uploader("Choose files to process", accept_multiple_files=True,
                type=['txt', 'pdf', 'docx', 'doc', 'mp3', 'mp4', 'wav', 'avi', 'mov', 'jpg', 'jpeg', 'png']
            ) # File object including file name no path stored in memory (streamlit run env)
            if uploaded_files:
                for uploaded_file in uploaded_files: 
                    file_metadata = {
                        "name": uploaded_file.name,
                        "type": uploaded_file.type,
                        "size": uploaded_file.size
                        }
        else:  # Google Drive
            if st.button("🔄**Refresh Browser**", help =  'Click to refresh Google Drive browser for updates', key="refresh_browser"): 
                # st.session_state.drive_folder_history = ['root']
                # st.session_state.drive_folder_names = {'root': 'GOOGLE DRIVE'}
                st.session_state.current_folder_files = None
                st.session_state.selected_files = []
                st.session_state.delete_confirm_needed = False
                st.session_state.overwrite_confirm_needed = False
                st.session_state.save_action = None
            if st.button("➡️📂**Main Folder for AI-Genenrated Stories**", key="goto_story_root_folder", help="Go to the AI-generated story folder"):
                    st.session_state.drive_folder_history = [
                        {'id': 'root', 'name': 'GOOGLE DRIVE'},
                        {'id': st.session_state.source_story_root_folder['id'], 'name': st.session_state.source_story_root_folder['name']},                     
                        {'id': st.session_state.generated_story_root_folder['id'], 'name': st.session_state.generated_story_root_folder['name']}]
                    st.session_state.current_folder_files = None
                    st.session_state.selected_files = []
                    st.session_state.delete_confirm_needed = False
                    st.session_state.overwrite_confirm_needed = False
                    st.session_state.save_action = None
                    st.rerun()                
            google_drive_browser()
        #st.header("")

        st.header("PROCESS HISTORY")
        for message in st.session_state.progress_history:
            st.text(message)        
        # Progress Bar
        if st.session_state.progress_value > 0:
            st.progress(st.session_state.progress_value / 100)


    # Main content area - show story tabs
    if st.session_state.story_content: 
        tabs = st.tabs(["👁️**GENERATED STORY VIEWER**", "🤖📝**AI STORY EDITOR**", "📝**MANUAL STORY EDITOR**"])
        with tabs[0]:
            col_html, col_docx = st.columns(2)
            with col_html:
                st.download_button("⬇️ DOWNLOAD STORY IN HTML FORMAT", data=story_html_format(),
                    file_name="story.html", mime="text/html")
            with col_docx:
                doc = docx.Document()
                markdown_to_docx(st.session_state.story_content, doc)
                docx_buffer = io.BytesIO()
                doc.save(docx_buffer)
                docx_buffer.seek(0)
                st.download_button("⬇️ DOWNLOAD STORY IN DOCX FORMAT", data=docx_buffer.getvalue(),
                    file_name="story.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
            # Inject custom CSS to match HTML font in Story Viewer
            st.markdown("""
                <style>
                    .story-viewer-body {
                        font-family: Arial, sans-serif !important;
                        font-size: 10pt !important;
                        margin: 10px !important;
                        text-align: justify !important;
                        padding: 0 !important;
                    }
                    .story-viewer-body h1 {
                        font-family: Arial, sans-serif !important;
                        font-size: 15pt !important;
                        text-align: center !important;
                        font-weight: bold !important;
                    }
                    .story-viewer-body h2 {
                        font-family: Arial, sans-serif !important;
                        font-size: 12pt !important;
                        text-align: justify !important;
                        font-weight: bold !important;
                    }
                    .story-viewer-body h3 {
                        font-family: Arial, sans-serif !important;
                        font-size: 11pt !important;
                        text-align: justify !important;
                        font-weight: bold !important;
                    }
                </style>
            """, unsafe_allow_html=True)
            # Wrap the story HTML in a div for CSS targeting
            st.markdown(f'<div class="story-viewer-body">{story_html_format()}</div>', unsafe_allow_html=True)

        with tabs[1]:
            st.text_area(
                "**ENTER YOUR INSTRUCTION TO REVISE THE CURRENT STORY WITH AI**:",
                height=300,
                key="story_revision_instruction",
                value=st.session_state.get("story_revision_instruction", ""))            
            col1, col2 = st.columns(2)
            with col1:
                if st.button("🔄**REVISE STORY**📝", key="revise_story_btn"):
                    if not st.session_state.story_content:
                        st.warning("No story content to revise!")
                    else:
                        with st.spinner("Revising story with AI is in progress ..."):
                            try:
                                revise_story()
                            except Exception as e:
                                st.error(f"Error generating story: {str(e)}")
            with col2: 
                if st.button("💾**SAVE REVISED STORY**", help ="Click to save revised story"):        
                    try: 
                        save_story_to_google_drive() 
                    except Exception as e:
                        st.error(f"Error saving revised story: {str(e)}") 
        with tabs[2]:
            soup = BeautifulSoup(story_html_format(), "html.parser")
            edited_texts = {} #store edited text
            tag_counter = {} # Track position-based identifiers for each tag
            for tag in soup.find_all():
                if tag.string and tag.name not in ["style"] and tag.parent.name == "body":
                    tag_name = tag.name
                    # Generate unique key per tag type using a counter
                    tag_counter[tag_name] = tag_counter.get(tag_name, 0) + 1
                    key = f"{tag_name}_{tag_counter[tag_name]}"
                    if tag_name in ["h2", "h3"]: # Display label instead of edit box for <h3>
                        st.write(f"{tag.string}")
                    else:
                        label = "Story Title" if tag_name == "h1" else " "   
                        edited_texts[key] = st.text_area(label=label, value=tag.string, height=100, key=f"textarea_{tag_name}_{tag_counter[tag_name]}")            
            if st.button("💾**SAVE CHANGES**", help ="Click to save edited changes to story"):        
                    try:
                        tag_counter = {}  # Reset for correct assignment
                        for tag in soup.find_all():
                            if tag.string and tag.name not in ["style"]: # not in ["style", "h3"]
                                tag_counter[tag.name] = tag_counter.get(tag.name, 0) + 1
                                key = f"{tag.name}_{tag_counter[tag.name]}"
                                if key in edited_texts:
                                    tag.string.replace_with(edited_texts[key])
                        st.session_state.story_content = html2text.html2text(str(soup)) 
                        save_story_to_google_drive() 
                    except Exception as e:
                        st.error(f"Error saving changes: {str(e)}")
if __name__ == "__main__":
    main()
